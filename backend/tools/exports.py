"""Export tools — turn structured content into real files on disk.

All tools share the same conventions:
- Files land at `data/exports/{turn_id}/{filename}` (turn_id comes from
  the runtime_events contextvar; "untagged" if missing).
- Filenames are sanitized — extensions are forced to match the tool kind.
- Tools return `(content_str, artifact_dict)` for LangChain's
  `content_and_artifact` response format. The string goes back to the
  model so it can reference the path; the artifact dict carries
  structured info for downstream consumers (frontend, logs, etc).
- Inputs are JSON strings (mirrors `delegate_to_departments`) so the
  schema is easy to render in tool calling and the model doesn't need
  to learn complex pydantic shapes.

Failures raise `ValueError` so LangChain's `ToolErrorMiddleware` can
surface the message to orc, which then retries with a corrected payload.
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Annotated, Any

from langchain_core.tools import tool

from config.paths import get_data_dir
from runtime_events import current_turn_id

logger = logging.getLogger(__name__)


_SAFE_FILENAME = re.compile(r"[^\w\-. 一-鿿]")


def _resolve_path(filename: str, ext: str) -> Path:
    """Map a requested filename to an absolute path under data/exports/{turn_id}/."""
    if not filename or not isinstance(filename, str):
        raise ValueError("filename must be a non-empty string")
    cleaned = _SAFE_FILENAME.sub("_", filename.strip()).strip("._-") or "output"
    # Force the requested extension (replacing any provided one).
    base = cleaned
    if "." in base:
        base = base.rsplit(".", 1)[0]
    final_name = f"{base}.{ext.lstrip('.')}"

    turn_id = current_turn_id() or "untagged"
    safe_turn = _SAFE_FILENAME.sub("_", str(turn_id)) or "untagged"
    out_dir = get_data_dir() / "exports" / safe_turn
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir / final_name


def _parse_json_payload(name: str, raw: str) -> Any:
    if not isinstance(raw, str):
        raise ValueError(f"{name} must be a JSON string")
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"{name} is not valid JSON: {exc}") from exc


def _artifact(path: Path, kind: str, extra: dict | None = None) -> dict:
    size = path.stat().st_size if path.exists() else 0
    out = {
        "kind": kind,
        "path": str(path),
        "filename": path.name,
        "size_bytes": size,
    }
    if extra:
        out.update(extra)
    return out


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------

@tool(response_format="content_and_artifact")
def create_pptx(
    filename: Annotated[
        str,
        "Desired filename without path. Example: '松林漫步-发布提案'. The .pptx extension is forced.",
    ],
    slides_json: Annotated[
        str,
        (
            "JSON array of slide objects. Each slide: {title: str, bullets: list[str], "
            "notes?: str, visual_hint?: str}. The first slide is treated as a title "
            "slide if its bullets list is empty. Aim for 5-15 slides."
        ),
    ],
) -> tuple[str, dict]:
    """Render a PowerPoint deck (.pptx) to disk.

    Use when the user explicitly asks for slides / a PPT / 幻灯片 / 发布提案,
    or when the deliverable is best consumed as a presentation. Each slide
    should be self-contained — content stays compact, with speaker notes
    holding the longer-form narrative.
    """
    from pptx import Presentation  # local import: avoid load cost on cold paths
    from pptx.util import Inches, Pt

    slides_data = _parse_json_payload("slides_json", slides_json)
    if not isinstance(slides_data, list) or not slides_data:
        raise ValueError("slides_json must be a non-empty JSON array")

    path = _resolve_path(filename, "pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]  # Title Slide
    content_layout = prs.slide_layouts[1]  # Title + Content

    for i, slide_obj in enumerate(slides_data):
        if not isinstance(slide_obj, dict):
            raise ValueError(f"slides_json[{i}] must be an object")
        title = str(slide_obj.get("title", "")).strip() or f"Slide {i + 1}"
        bullets = slide_obj.get("bullets") or []
        if not isinstance(bullets, list):
            raise ValueError(f"slides_json[{i}].bullets must be a list of strings")
        notes = str(slide_obj.get("notes", "") or "").strip()
        visual_hint = str(slide_obj.get("visual_hint", "") or "").strip()

        is_cover = i == 0 and not bullets
        slide = prs.slides.add_slide(title_layout if is_cover else content_layout)
        slide.shapes.title.text = title

        if is_cover:
            # Subtitle placeholder for title slides
            if len(slide.placeholders) > 1 and visual_hint:
                slide.placeholders[1].text = visual_hint
        else:
            body = slide.placeholders[1].text_frame
            body.clear()
            for j, b in enumerate(bullets):
                text = str(b).strip()
                if not text:
                    continue
                p = body.paragraphs[0] if j == 0 else body.add_paragraph()
                p.text = text
                p.font.size = Pt(20)

        # Notes — keep speaker-notes-style narrative + visual hints together
        combined_notes = "\n\n".join(part for part in (notes, f"视觉建议：{visual_hint}" if visual_hint else "") if part)
        if combined_notes:
            slide.notes_slide.notes_text_frame.text = combined_notes

    prs.save(str(path))
    artifact = _artifact(path, "pptx", {"slide_count": len(slides_data)})
    msg = f"Wrote {len(slides_data)}-slide PPT to {path} ({artifact['size_bytes']} bytes)."
    return msg, artifact


# ---------------------------------------------------------------------------
# .docx
# ---------------------------------------------------------------------------

@tool(response_format="content_and_artifact")
def create_docx(
    filename: Annotated[
        str,
        "Desired filename without path. Example: '经销商合同要点'. The .docx extension is forced.",
    ],
    title: Annotated[str, "Document title shown at the top (level-0 heading)."],
    sections_json: Annotated[
        str,
        (
            "JSON array of section objects. Each section: "
            "{heading: str, level?: int (default 1, 1-3), paragraphs: list[str]}. "
            "Paragraphs are rendered as separate <p> blocks."
        ),
    ],
) -> tuple[str, dict]:
    """Render a Word document (.docx) to disk.

    Use for long-form deliverables: 合同要点 / 公关稿 / SOP / 培训材料 / 风险评估报告.
    Prefer this over PPT when the content is paragraph-heavy rather than bullet-heavy.
    """
    from docx import Document

    sections = _parse_json_payload("sections_json", sections_json)
    if not isinstance(sections, list):
        raise ValueError("sections_json must be a JSON array")

    path = _resolve_path(filename, "docx")
    doc = Document()
    doc.add_heading(str(title or "Document"), level=0)

    for i, sec in enumerate(sections):
        if not isinstance(sec, dict):
            raise ValueError(f"sections_json[{i}] must be an object")
        heading = str(sec.get("heading", "") or "").strip()
        level = int(sec.get("level", 1) or 1)
        level = max(1, min(3, level))
        if heading:
            doc.add_heading(heading, level=level)
        paragraphs = sec.get("paragraphs") or []
        if not isinstance(paragraphs, list):
            raise ValueError(f"sections_json[{i}].paragraphs must be a list of strings")
        for p in paragraphs:
            text = str(p).strip()
            if text:
                doc.add_paragraph(text)

    doc.save(str(path))
    artifact = _artifact(path, "docx", {"section_count": len(sections)})
    msg = f"Wrote Word doc with {len(sections)} section(s) to {path} ({artifact['size_bytes']} bytes)."
    return msg, artifact


# ---------------------------------------------------------------------------
# .xlsx
# ---------------------------------------------------------------------------

@tool(response_format="content_and_artifact")
def create_xlsx(
    filename: Annotated[
        str,
        "Desired filename without path. Example: '秋冬OTB-货盘表'. The .xlsx extension is forced.",
    ],
    sheets_json: Annotated[
        str,
        (
            "JSON array of sheet objects. Each sheet: "
            "{name: str, headers: list[str], rows: list[list[Any]]}. "
            "Cell values can be strings, numbers, booleans, or null. "
            "Aim for short tables (< 1000 rows) — this is a planning artifact, not a data dump."
        ),
    ],
) -> tuple[str, dict]:
    """Render an Excel workbook (.xlsx) to disk.

    Use for tabular deliverables: OTB / 货盘表 / 财务测算 / 售罄分析 / 库存盘点 /
    营销日历 / 投放排期. Each sheet becomes a tab in the workbook. Headers are
    bolded automatically.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font

    sheets = _parse_json_payload("sheets_json", sheets_json)
    if not isinstance(sheets, list) or not sheets:
        raise ValueError("sheets_json must be a non-empty JSON array")

    path = _resolve_path(filename, "xlsx")
    wb = Workbook()
    wb.remove(wb.active)  # Drop default blank sheet
    bold = Font(bold=True)
    total_rows = 0

    for i, sheet in enumerate(sheets):
        if not isinstance(sheet, dict):
            raise ValueError(f"sheets_json[{i}] must be an object")
        name = str(sheet.get("name", "") or f"Sheet{i + 1}")[:31]  # Excel max
        headers = sheet.get("headers") or []
        rows = sheet.get("rows") or []
        if not isinstance(headers, list) or not isinstance(rows, list):
            raise ValueError(f"sheets_json[{i}]: headers and rows must be lists")
        ws = wb.create_sheet(title=name)
        if headers:
            ws.append([str(h) for h in headers])
            for cell in ws[1]:
                cell.font = bold
        for row in rows:
            if not isinstance(row, list):
                raise ValueError(f"sheets_json[{i}]: each row must be a list")
            ws.append([
                v if isinstance(v, (int, float, bool)) or v is None else str(v)
                for v in row
            ])
        total_rows += len(rows)

    wb.save(str(path))
    artifact = _artifact(path, "xlsx", {"sheet_count": len(sheets), "row_count": total_rows})
    msg = f"Wrote Excel workbook with {len(sheets)} sheet(s), {total_rows} row(s) to {path} ({artifact['size_bytes']} bytes)."
    return msg, artifact


# ---------------------------------------------------------------------------
# .md
# ---------------------------------------------------------------------------

@tool(response_format="content_and_artifact")
def create_markdown(
    filename: Annotated[
        str,
        "Desired filename without path. Example: '调研笔记-松林漫步'. The .md extension is forced.",
    ],
    content: Annotated[
        str,
        "Full markdown content. Use standard markdown (headings, lists, tables, code blocks).",
    ],
) -> tuple[str, dict]:
    """Render a Markdown report (.md) to disk.

    Use for lightweight long-form text: 备忘录 / 调研笔记 / 复盘 / 内部周报 /
    简短分析报告. Cheap and always succeeds — pick this when a .docx feels too
    heavyweight for the deliverable.
    """
    if not isinstance(content, str) or not content.strip():
        raise ValueError("content must be a non-empty string")
    path = _resolve_path(filename, "md")
    path.write_text(content, encoding="utf-8")
    artifact = _artifact(path, "md", {"chars": len(content)})
    msg = f"Wrote markdown ({len(content)} chars) to {path} ({artifact['size_bytes']} bytes)."
    return msg, artifact
