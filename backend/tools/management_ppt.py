"""Management-grade PPT builder — replaces the bare `create_pptx` tool.

Reads a structured deck spec (JSON) and renders a deck that:
- Carries a brand theme (palette, fonts, footer)
- Has six slide *types* with distinct visual treatments:
    cover     — full-color background, oversized title
    agenda    — numbered chapter list
    divider   — full-color section break with chapter title
    content   — title + 1-N content blocks (bullets, cards, callouts)
    data      — title + matplotlib chart + insight callout
    closing   — action items + decision points

The builder is intentionally opinionated: layouts are hand-positioned
(not from a .pptx template) so the output is deterministic and the
visual quality doesn't depend on user-supplied master slides.
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Annotated, Any

from langchain_core.tools import tool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from runtime_events import current_turn_id
from tools.charts import make_bar_chart, make_line_chart, make_pie_chart
from tools.exports import _resolve_path, _artifact, _parse_json_payload
from tools.ppt_theme import PPTTheme, get_theme

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Geometry helpers
# ---------------------------------------------------------------------------

SLIDE_W_IN = 13.333  # 16:9 widescreen
SLIDE_H_IN = 7.5


def _set_font(run, *, size: int, color: RGBColor, bold: bool = False, theme: PPTTheme) -> None:
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = theme.font_latin
    # Set east-asian font via XML — python-pptx doesn't expose this directly
    rPr = run._r.get_or_add_rPr()
    # Remove existing eastAsia tag if any (for re-runs)
    for tag in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}ea"):
        rPr.remove(tag)
    from lxml import etree
    ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", theme.font_east_asian)


def _add_text(slide, *, x: float, y: float, w: float, h: float,
              text: str, theme: PPTTheme, size: int, color: RGBColor | None = None,
              bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT,
              vertical_center: bool = False) -> Any:
    """Add a textbox with theme-aware font + a single paragraph."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if vertical_center:
        tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE = 3
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=color or theme.text, bold=bold, theme=theme)
    return box


def _add_paragraphs(slide, *, x: float, y: float, w: float, h: float,
                    lines: list[str], theme: PPTTheme, size: int = 16,
                    color: RGBColor | None = None, bullet: bool = True,
                    space_after_pt: int = 6) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after_pt)
        if bullet:
            line = f"•  {line}"
        run = p.add_run()
        run.text = line
        _set_font(run, size=size, color=color or theme.text, theme=theme)
    return box


def _add_rect(slide, *, x: float, y: float, w: float, h: float,
              fill: RGBColor, line: RGBColor | None = None) -> Any:
    """Solid rectangle (no shadow). Used for color blocks and cards."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_footer(slide, *, theme: PPTTheme, page_num: int | None, brand: str, total: int | None = None) -> None:
    # Left: brand name
    _add_text(slide, x=0.6, y=SLIDE_H_IN - 0.32, w=6.0, h=0.25,
              text=brand, theme=theme, size=theme.footer, color=theme.text_muted)
    # Right: page number
    if page_num is not None:
        page_text = f"{page_num} / {total}" if total else f"{page_num}"
        _add_text(slide, x=SLIDE_W_IN - 1.6, y=SLIDE_H_IN - 0.32, w=1.0, h=0.25,
                  text=page_text, theme=theme, size=theme.footer, color=theme.text_muted,
                  align=PP_ALIGN.RIGHT)


def _add_title_bar(slide, *, title: str, theme: PPTTheme, eyebrow: str | None = None) -> None:
    """Title bar at top: optional eyebrow (small uppercase) + main title +
    thin accent rule beneath. Used by content / data / closing slides."""
    y = theme.margin_top
    if eyebrow:
        _add_text(slide, x=theme.margin_left, y=y, w=SLIDE_W_IN - 2, h=0.25,
                  text=eyebrow.upper(), theme=theme, size=10, color=theme.accent, bold=True)
        y += 0.25
    _add_text(slide, x=theme.margin_left, y=y, w=SLIDE_W_IN - 2 * theme.margin_left, h=theme.title_height,
              text=title, theme=theme, size=theme.title_slide, color=theme.text, bold=True)
    # accent rule
    rule_y = y + theme.title_height - 0.1
    _add_rect(slide, x=theme.margin_left, y=rule_y, w=0.6, h=0.04, fill=theme.accent)


def _normalize_chart_series(raw: Any) -> dict[str, list[float]]:
    """Accept both {"name": [values]} and [{"name": ..., "values": ...}]."""
    if isinstance(raw, dict):
        out: dict[str, list[float]] = {}
        for name, values in raw.items():
            if isinstance(values, list):
                out[str(name)] = [float(v) for v in values]
        return out
    if isinstance(raw, list):
        out: dict[str, list[float]] = {}
        for i, item in enumerate(raw, 1):
            if not isinstance(item, dict):
                continue
            values = item.get("values", [])
            if isinstance(values, list):
                out[str(item.get("name") or f"Series {i}")] = [float(v) for v in values]
        return out
    return {}


def _normalize_actions(raw: Any) -> list[dict[str, str]]:
    """Accept action dicts or plain action strings from LLM-generated specs."""
    if isinstance(raw, dict):
        raw = [raw]
    if not isinstance(raw, list):
        return []
    out: list[dict[str, str]] = []
    for item in raw:
        if isinstance(item, dict):
            out.append({
                "owner": str(item.get("owner", "")),
                "action": str(item.get("action", "")),
                "due": str(item.get("due", "")),
            })
        else:
            action = str(item).strip()
            if action:
                out.append({"owner": "", "action": action, "due": ""})
    return out


def _normalize_decisions(raw: Any) -> list[str]:
    """Accept decision strings or {"question": ...} objects."""
    if isinstance(raw, list):
        out = []
        for item in raw:
            value = item.get("question", "") if isinstance(item, dict) else item
            text = str(value).strip()
            if text:
                out.append(text)
        return out
    text = str(raw or "").strip()
    return [text] if text else []


# ---------------------------------------------------------------------------
# Slide renderers
# ---------------------------------------------------------------------------

def _render_cover(slide, spec: dict, theme: PPTTheme) -> None:
    """Full-color cover slide.

    spec keys: title, subtitle (optional), date (optional), author (optional)
    """
    _set_slide_bg(slide, theme.primary)
    # Accent vertical bar on left
    _add_rect(slide, x=0, y=0, w=0.35, h=SLIDE_H_IN, fill=theme.accent)
    # Eyebrow
    eyebrow = spec.get("eyebrow", "boxcc 协作汇报").upper()
    _add_text(slide, x=1.2, y=2.4, w=10, h=0.4, text=eyebrow,
              theme=theme, size=12, color=theme.accent, bold=True)
    # Title
    _add_text(slide, x=1.2, y=2.9, w=11, h=2,
              text=str(spec.get("title", "Untitled Deck")),
              theme=theme, size=theme.title_cover, color=theme.bg, bold=True)
    # Subtitle
    subtitle = spec.get("subtitle")
    if subtitle:
        _add_text(slide, x=1.2, y=4.6, w=11, h=0.7,
                  text=str(subtitle), theme=theme, size=18, color=theme.subtle)
    # Footer line: author / date
    meta_parts = []
    if spec.get("author"):
        meta_parts.append(str(spec["author"]))
    if spec.get("date"):
        meta_parts.append(str(spec["date"]))
    if meta_parts:
        _add_text(slide, x=1.2, y=SLIDE_H_IN - 0.9, w=11, h=0.3,
                  text=" · ".join(meta_parts), theme=theme, size=11, color=theme.subtle)


def _render_agenda(slide, spec: dict, theme: PPTTheme) -> None:
    """Agenda / table of contents."""
    _set_slide_bg(slide, theme.bg)
    _add_title_bar(slide, title=spec.get("title", "本次议程"), theme=theme, eyebrow="agenda")
    items: list[str] = spec.get("items", [])
    y = 2.0
    for i, item in enumerate(items, 1):
        # number circle
        _add_rect(slide, x=theme.margin_left, y=y, w=0.5, h=0.5,
                  fill=theme.subtle)
        _add_text(slide, x=theme.margin_left, y=y, w=0.5, h=0.5,
                  text=f"{i:02d}", theme=theme, size=18, color=theme.primary,
                  bold=True, align=PP_ALIGN.CENTER, vertical_center=True)
        _add_text(slide, x=theme.margin_left + 0.85, y=y + 0.06, w=SLIDE_W_IN - 2.5, h=0.5,
                  text=str(item), theme=theme, size=18, color=theme.text)
        y += 0.72
        if y > SLIDE_H_IN - 1:
            break


def _render_divider(slide, spec: dict, theme: PPTTheme) -> None:
    """Section divider — full-color page with chapter number + title."""
    _set_slide_bg(slide, theme.primary)
    # Big number
    num = str(spec.get("number", ""))
    if num:
        _add_text(slide, x=1.0, y=2.0, w=4, h=2,
                  text=num, theme=theme, size=96, color=theme.accent, bold=True)
    # Chapter title
    _add_text(slide, x=1.0, y=4.2, w=11, h=1.5,
              text=str(spec.get("title", "")), theme=theme, size=theme.title_section,
              color=theme.bg, bold=True)
    # Optional intro line
    intro = spec.get("intro")
    if intro:
        _add_text(slide, x=1.0, y=5.7, w=11, h=1,
                  text=str(intro), theme=theme, size=16, color=theme.subtle)


def _render_content(slide, spec: dict, theme: PPTTheme) -> None:
    """Content slide.

    spec keys:
      title: str
      eyebrow: str (optional)
      lead: str (optional, one-sentence takeaway at top of body)
      blocks: list of one of:
        {"type": "bullets", "items": [str, ...]}
        {"type": "cards", "items": [{"title": str, "desc": str}, ...]}  # up to 4
        {"type": "callout", "text": str}
        {"type": "paragraph", "text": str}
        {"type": "kv", "items": [{"label": str, "value": str}, ...]}
    """
    _set_slide_bg(slide, theme.bg)
    _add_title_bar(slide, title=spec.get("title", ""), theme=theme, eyebrow=spec.get("eyebrow"))
    body_top = theme.margin_top + theme.title_height + 0.4
    if spec.get("lead"):
        _add_text(slide, x=theme.margin_left, y=body_top, w=SLIDE_W_IN - 2 * theme.margin_left, h=0.6,
                  text=str(spec["lead"]), theme=theme, size=18, color=theme.primary, bold=True)
        body_top += 0.7

    body_left = theme.margin_left
    body_w = SLIDE_W_IN - 2 * theme.margin_left
    cursor = body_top

    for block in spec.get("blocks", []):
        btype = block.get("type", "bullets")
        if btype == "bullets":
            items = [str(x) for x in block.get("items", [])]
            h = max(0.3, 0.35 * len(items)) + 0.1
            _add_paragraphs(slide, x=body_left, y=cursor, w=body_w, h=h,
                            lines=items, theme=theme, size=theme.body, bullet=True)
            cursor += h + 0.15
        elif btype == "paragraph":
            text = str(block.get("text", ""))
            h = max(0.4, len(text) / 90 * 0.3)
            _add_text(slide, x=body_left, y=cursor, w=body_w, h=h,
                      text=text, theme=theme, size=theme.body)
            cursor += h + 0.15
        elif btype == "callout":
            text = str(block.get("text", ""))
            h = max(0.6, len(text) / 80 * 0.3) + 0.2
            _add_rect(slide, x=body_left, y=cursor, w=body_w, h=h, fill=theme.subtle)
            # Left accent stripe
            _add_rect(slide, x=body_left, y=cursor, w=0.08, h=h, fill=theme.accent)
            _add_text(slide, x=body_left + 0.3, y=cursor + 0.15, w=body_w - 0.5, h=h - 0.2,
                      text=text, theme=theme, size=theme.callout, color=theme.text, bold=False)
            cursor += h + 0.2
        elif btype == "cards":
            items = block.get("items", [])[:4]
            n = max(1, len(items))
            gap = 0.2
            card_w = (body_w - gap * (n - 1)) / n
            card_h = 1.7
            for i, it in enumerate(items):
                cx = body_left + i * (card_w + gap)
                _add_rect(slide, x=cx, y=cursor, w=card_w, h=card_h, fill=theme.subtle)
                # Index dot
                _add_text(slide, x=cx + 0.15, y=cursor + 0.12, w=0.4, h=0.3,
                          text=f"{i+1:02d}", theme=theme, size=12, color=theme.accent, bold=True)
                _add_text(slide, x=cx + 0.15, y=cursor + 0.42, w=card_w - 0.3, h=0.45,
                          text=str(it.get("title", "")), theme=theme, size=15,
                          color=theme.text, bold=True)
                _add_text(slide, x=cx + 0.15, y=cursor + 0.85, w=card_w - 0.3, h=card_h - 0.95,
                          text=str(it.get("desc", "")), theme=theme, size=11, color=theme.text_muted)
            cursor += card_h + 0.2
        elif btype == "kv":
            items = block.get("items", [])
            for kv in items:
                label = str(kv.get("label", ""))
                value = str(kv.get("value", ""))
                _add_text(slide, x=body_left, y=cursor, w=2.6, h=0.35,
                          text=label, theme=theme, size=13, color=theme.text_muted, bold=False)
                _add_text(slide, x=body_left + 2.7, y=cursor, w=body_w - 2.7, h=0.35,
                          text=value, theme=theme, size=14, color=theme.text)
                cursor += 0.4
            cursor += 0.1


def _render_data(slide, spec: dict, theme: PPTTheme) -> None:
    """Data slide: title + chart + takeaway sentence.

    spec keys:
      title: str
      takeaway: str (the single sentence above the chart)
      chart: {"type": "bar"|"line"|"pie", ...payload...}
    """
    _set_slide_bg(slide, theme.bg)
    _add_title_bar(slide, title=spec.get("title", ""), theme=theme, eyebrow=spec.get("eyebrow", "data"))
    body_top = theme.margin_top + theme.title_height + 0.4
    if spec.get("takeaway"):
        _add_text(slide, x=theme.margin_left, y=body_top, w=SLIDE_W_IN - 2 * theme.margin_left, h=0.6,
                  text=str(spec["takeaway"]), theme=theme, size=18, color=theme.primary, bold=True)
        body_top += 0.7

    chart = spec.get("chart") or {}
    ctype = chart.get("type", "bar")
    try:
        series = _normalize_chart_series(chart.get("series", {}))
        if ctype == "line":
            png = make_line_chart(theme=theme,
                                  categories=chart.get("categories", []),
                                  series=series,
                                  title=chart.get("title", ""))
        elif ctype == "pie":
            png = make_pie_chart(theme=theme,
                                 labels=chart.get("labels", []),
                                 values=chart.get("values", []),
                                 title=chart.get("title", ""))
        else:
            png = make_bar_chart(theme=theme,
                                 categories=chart.get("categories", []),
                                 series=series,
                                 title=chart.get("title", ""))
        chart_h = SLIDE_H_IN - body_top - 1.0
        chart_w = SLIDE_W_IN - 2 * theme.margin_left
        slide.shapes.add_picture(str(png), Inches(theme.margin_left), Inches(body_top),
                                 width=Inches(chart_w), height=Inches(chart_h))
    except Exception as exc:
        _add_text(slide, x=theme.margin_left, y=body_top, w=10, h=1,
                  text=f"(chart render failed: {exc})", theme=theme, size=14, color=theme.text_muted)


def _render_closing(slide, spec: dict, theme: PPTTheme) -> None:
    """Closing slide: action items + decision questions.

    spec keys:
      title: str (default "下一步与决策")
      actions: list[{"owner": str, "action": str, "due": str}]
      decisions: list[str]
    """
    _set_slide_bg(slide, theme.bg)
    _add_title_bar(slide, title=spec.get("title", "下一步与决策"), theme=theme, eyebrow="next steps")
    body_top = theme.margin_top + theme.title_height + 0.4

    actions = _normalize_actions(spec.get("actions", []))
    decisions = _normalize_decisions(spec.get("decisions", []))

    # Two-column layout
    col_w = (SLIDE_W_IN - 2 * theme.margin_left - 0.4) / 2

    # Left column: action items
    if actions:
        _add_text(slide, x=theme.margin_left, y=body_top, w=col_w, h=0.4,
                  text="行动项".upper(), theme=theme, size=11, color=theme.accent, bold=True)
        y = body_top + 0.5
        for a in actions[:8]:
            owner = str(a.get("owner", ""))
            action = str(a.get("action", ""))
            due = str(a.get("due", ""))
            _add_rect(slide, x=theme.margin_left, y=y, w=col_w, h=0.7, fill=theme.subtle)
            _add_text(slide, x=theme.margin_left + 0.2, y=y + 0.08, w=col_w - 0.4, h=0.28,
                      text=action, theme=theme, size=13, color=theme.text, bold=True)
            meta = " · ".join(filter(None, [owner, f"截止：{due}" if due else None]))
            _add_text(slide, x=theme.margin_left + 0.2, y=y + 0.38, w=col_w - 0.4, h=0.25,
                      text=meta, theme=theme, size=10, color=theme.text_muted)
            y += 0.8
            if y > SLIDE_H_IN - 1:
                break

    # Right column: decision questions
    if decisions:
        rx = theme.margin_left + col_w + 0.4
        _add_text(slide, x=rx, y=body_top, w=col_w, h=0.4,
                  text="待决策".upper(), theme=theme, size=11, color=theme.accent, bold=True)
        y = body_top + 0.5
        for d in decisions[:6]:
            _add_text(slide, x=rx, y=y, w=col_w, h=0.7,
                      text=f"?  {d}", theme=theme, size=14, color=theme.text)
            y += 0.7
            if y > SLIDE_H_IN - 1:
                break


# ---------------------------------------------------------------------------
# Top-level builder
# ---------------------------------------------------------------------------

SLIDE_RENDERERS = {
    "cover": _render_cover,
    "agenda": _render_agenda,
    "divider": _render_divider,
    "content": _render_content,
    "data": _render_data,
    "closing": _render_closing,
}


@tool(response_format="content_and_artifact")
def create_management_ppt(
    filename: Annotated[
        str,
        "Desired filename without path. Example: '松林漫步-季度复盘'. The .pptx extension is forced.",
    ],
    deck_json: Annotated[
        str,
        (
            "JSON object describing the deck. Top-level keys: "
            "meta: {title, subtitle?, author?, date?, brand?, theme?}; "
            "slides: list of slide objects each {type: cover|agenda|divider|content|data|closing, ...payload}. "
            "Theme defaults to 'editorial'. Brand text shows in footer. "
            "Use 'data' slides for chart-driven pages; chart payload is {type:'bar'|'line'|'pie', categories?, series?|labels?, values?, title?}."
        ),
    ],
) -> tuple[str, dict]:
    """Render a management-grade PPT deck with theming, layouts, and charts.

    Use this instead of the bare `create_pptx` when the deliverable is a
    management briefing / quarterly review / launch proposal / strategy
    deck — anything that needs to look like it came out of a brand team,
    not a typewriter.

    Structure your deck like a magazine: cover -> agenda -> chapter
    dividers -> content slides -> data slides -> closing with actions
    and decisions. Keep each slide focused on one idea.
    """
    deck = _parse_json_payload("deck_json", deck_json)
    if not isinstance(deck, dict):
        raise ValueError("deck_json must be a JSON object")

    meta = deck.get("meta", {})
    if not isinstance(meta, dict):
        meta = {}
    slides = deck.get("slides", [])
    if not isinstance(slides, list) or not slides:
        raise ValueError("deck_json.slides must be a non-empty list")

    theme = get_theme(meta.get("theme"))
    brand = str(meta.get("brand") or "boxcc")

    path = _resolve_path(filename, "pptx")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]  # blank — we hand-position everything

    total = len(slides)
    chapter_counter = 0
    for i, sd in enumerate(slides, 1):
        if not isinstance(sd, dict):
            raise ValueError(f"slides[{i-1}] must be an object")
        stype = sd.get("type", "content")
        renderer = SLIDE_RENDERERS.get(stype, _render_content)
        slide = prs.slides.add_slide(blank_layout)

        # Auto-number dividers if no explicit number provided
        if stype == "divider" and not sd.get("number"):
            chapter_counter += 1
            sd = {**sd, "number": f"{chapter_counter:02d}"}

        renderer(slide, sd, theme)

        # Add footer to every slide except cover (the cover already carries meta)
        if stype != "cover":
            _add_footer(slide, theme=theme, page_num=i, total=total, brand=brand)

        # Speaker notes pass-through
        notes = sd.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(str(path))
    artifact = _artifact(path, "pptx", {
        "slide_count": total,
        "theme": theme.name,
        "brand": brand,
        "kind_detail": "management_ppt",
    })
    msg = f"Wrote {total}-slide management PPT to {path} (theme={theme.name}, {artifact['size_bytes']} bytes)."
    return msg, artifact
